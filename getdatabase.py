import os
import glob
from typing import List, Tuple, Dict
import json
import numpy as np
from tqdm import tqdm
from sentence_transformers import SentenceTransformer
import docx2txt
import pdfplumber
from pptx import Presentation
import logging
from datetime import datetime
import re

# Langchain imports
from langchain.vectorstores import Chroma
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.schema import Document

# 设置日志
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# 设置全局变量
EMBEDDING_MODEL = 'paraphrase-multilingual-MiniLM-L12-v2'
CHUNK_SIZE = 1000
OVERLAP = 200

def setup_file_logger(output_dir):
    log_file = os.path.join(output_dir, f'processing_log_{datetime.now().strftime("%Y%m%d_%H%M%S")}.txt')
    file_handler = logging.FileHandler(log_file, encoding='utf-8')
    file_handler.setLevel(logging.INFO)
    file_handler.setFormatter(logging.Formatter('%(message)s'))

    file_logger = logging.getLogger('file_logger')
    file_logger.setLevel(logging.INFO)
    file_logger.addHandler(file_handler)
    file_logger.propagate = False

    return log_file, file_logger, file_handler

def log_processing_result(file_logger, file_path, success, message=''):
    status = "成功" if success else "失败"
    file_logger.info(f"{file_path}: {status} - {message}")


def clean_text(text: str) -> str:
    # 统一换行符
    text = text.replace('\r\n', '\n').replace('\r', '\n')
    
    # 移除连续的换行符，替换为单个换行符
    text = re.sub(r'\n{2,}', '\n', text)
    
    # 移除行首和行尾的空白字符
    lines = [line.strip() for line in text.split('\n')]
    
    # 处理目录样式的行
    cleaned_lines = []
    for line in lines:
        # 移除行末连续的点号和数字
        line = re.sub(r'\.{2,}\s*\d+\s*$', '', line)
        # 移除整行都是点号的情况
        if not re.match(r'^\.+$', line):
            # 移除行中间过多的连续点号（超过3个）
            line = re.sub(r'\.{3,}', '...', line)
            cleaned_lines.append(line)
    
    # 重新组合文本
    text = ' '.join(cleaned_lines)
    
    # 移除连续的空格，替换为单个空格
    text = re.sub(r'\s+', ' ', text)
    
    # 处理章节编号，确保编号和文本之间有一个空格
    text = re.sub(r'(\d+(\.\d+)*)', r' \1 ', text)
    text = re.sub(r'\s+', ' ', text)  # 再次清理可能产生的多余空格
    
    # 最后再次去除首尾空白
    text = text.strip()
    
    return text


def extract_text_from_file(file_path: str) -> str:
    file_extension = os.path.splitext(file_path)[1].lower()
    logging.info(f"Extracting text from file: {file_path}")

    try:
        if file_extension in ['.doc', '.docx']:
            text = docx2txt.process(file_path)
        elif file_extension == '.pdf':
            with pdfplumber.open(file_path) as pdf:
                text = "\n".join(page.extract_text() or "" for page in pdf.pages)
        elif file_extension == '.pptx':
            prs = Presentation(file_path)
            text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text'))
        elif file_extension in ['.txt', '.md']:
            encodings = ['utf-8', 'gbk', 'iso-8859-1']
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        text = f.read()
                    break
                except UnicodeDecodeError:
                    continue
            else:
                raise ValueError(f"Unable to decode file {file_path} with any of the tried encodings")
        else:
            raise ValueError(f"Unsupported file format: {file_extension}")

        # 应用增强的文本清理
        text = clean_text(text)

        return text
    except Exception as e:
        logging.error(f"Error extracting text from {file_path}: {str(e)}")
        return ""

def split_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = OVERLAP) -> List[str]:
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        chunks.append(chunk)
        start = end - overlap
    return chunks

def process_documents(directory: str, file_logger) -> Tuple[List[Dict[str, str]], np.ndarray]:
    logging.info(f"正在加载模型 {EMBEDDING_MODEL}...")
    model = SentenceTransformer(EMBEDDING_MODEL)
    
    all_chunks = []
    all_vectors = []
    
    total_files = sum([len(files) for r, d, files in os.walk(directory)])
    logging.info(f"Found {total_files} files in the directory and its subdirectories")

    with tqdm(total=total_files, desc="Processing files") as pbar:
        for root, dirs, files in os.walk(directory):
            for file in files:
                file_path = os.path.join(root, file)
                try:
                    file_extension = os.path.splitext(file_path)[1].lower()
                    if file_extension in ['.doc', '.docx', '.pdf', '.pptx']:
                        text = extract_text_from_file(file_path)
                        
                        if not text:
                            logging.warning(f"Skipping file due to empty content: {file_path}")
                            log_processing_result(file_logger, file_path, False, "文件内容为空，已跳过")
                            pbar.update(1)
                            continue

                        chunks = split_text(text)
                        vectors = model.encode(chunks)
                        
                        for chunk in chunks:
                            all_chunks.append({
                                "text": chunk,
                                "file_path": file_path
                            })
                        all_vectors.append(vectors)
                        logging.info(f"Processed file: {file_path}, chunks: {len(chunks)}")
                        log_processing_result(file_logger, file_path, True, f"成功处理，生成 {len(chunks)} 个文本块")
                    else:
                        logging.warning(f"Skipping unsupported file: {file_path}")
                        log_processing_result(file_logger, file_path, False, "不支持的文件格式")
                    
                except Exception as e:
                    logging.error(f"处理文件 {file_path} 时出错: {str(e)}")
                    log_processing_result(file_logger, file_path, False, f"处理出错: {str(e)}")
                
                pbar.update(1)
    
    if all_vectors:
        all_vectors = np.vstack(all_vectors)
        logging.info(f"Total vectors: {all_vectors.shape}")
    else:
        logging.warning("No vectors generated")
        all_vectors = np.array([])

    return all_chunks, all_vectors

def build_and_save_chroma_db(chunks: List[Dict[str, str]], vectors: np.ndarray, output_dir: str):
    if vectors.size == 0:
        logging.warning("No vectors to build database, skipping database creation")
        return

    logging.info("构建和保存Chroma数据库...")
    embeddings = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL)
    
    # 创建符合Langchain Document格式的文档列表
    documents = [
        Document(
            page_content=chunk["text"],
            metadata={"source": chunk["file_path"]}
        ) for chunk in chunks
    ]
    
    vectorstore = Chroma.from_documents(
        documents,
        embeddings,
        persist_directory=output_dir
    )
    vectorstore.persist()
    logging.info(f"Chroma数据库已保存至 {output_dir}")

def main():
    input_dir = "/root/trytry/kb/风电知识库"  #数据库文件地址
    output_dir = "./chroma_db"
    chunks_path = os.path.join(output_dir, "chunks.json")

    os.makedirs(output_dir, exist_ok=True)
    log_file, file_logger, file_handler = setup_file_logger(output_dir)

    logging.info("开始处理文档...")
    chunks, vectors = process_documents(input_dir, file_logger)
    
    if len(chunks) > 0 and vectors.size > 0:
        logging.info("构建和保存Chroma数据库...")
        build_and_save_chroma_db(chunks, vectors, output_dir)
    else:
        logging.warning("No chunks or vectors generated, skipping database creation")
    
    logging.info("保存文本块...")
    with open(chunks_path, 'w', encoding='utf-8') as f:
        json.dump(chunks, f, ensure_ascii=False, indent=2)
    
    logging.info(f"处理完成！生成的文件：\n1. {output_dir} (Chroma数据库目录)\n2. {chunks_path}\n3. {log_file}")

    # 关闭文件日志处理器
    file_handler.close()
    file_logger.removeHandler(file_handler)

if __name__ == "__main__":
    main() 